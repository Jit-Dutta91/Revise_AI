# backend/file_reader.py
# Extracts text from uploaded files (PDF, DOCX, PPTX, TXT, images)

import io


def extract_text_from_file(uploaded_file) -> str:
    """
    Extract text content from various file types.
    Supports: PDF, DOCX, PPTX, TXT, MD
    Images: returns a note (Groq text models can't process images directly)
    """

    filename = uploaded_file.name.lower()
    file_bytes = uploaded_file.read()

    # ── TXT / MD ──────────────────────────────────────────────
    if filename.endswith((".txt", ".md")):
        try:
            return file_bytes.decode("utf-8", errors="ignore")
        except Exception:
            return ""

    # ── PDF ───────────────────────────────────────────────────
    elif filename.endswith(".pdf"):
        try:
            import pypdf
            reader = pypdf.PdfReader(io.BytesIO(file_bytes))
            text = ""
            for page in reader.pages:
                text += page.extract_text() or ""
            return text.strip()
        except ImportError:
            return f"[PDF file: {uploaded_file.name} — install pypdf to extract text]"
        except Exception as e:
            return f"[Could not read PDF: {e}]"

    # ── DOCX ──────────────────────────────────────────────────
    elif filename.endswith((".docx", ".doc")):
        try:
            import docx
            doc = docx.Document(io.BytesIO(file_bytes))
            text = "\n".join([para.text for para in doc.paragraphs if para.text.strip()])
            return text.strip()
        except ImportError:
            return f"[DOCX file: {uploaded_file.name} — install python-docx to extract text]"
        except Exception as e:
            return f"[Could not read DOCX: {e}]"

    # ── PPTX ──────────────────────────────────────────────────
    elif filename.endswith((".pptx", ".ppt")):
        try:
            from pptx import Presentation
            prs  = Presentation(io.BytesIO(file_bytes))
            text = ""
            for slide_num, slide in enumerate(prs.slides, 1):
                text += f"\n[Slide {slide_num}]\n"
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        text += shape.text.strip() + "\n"
            return text.strip()
        except ImportError:
            return f"[PPTX file: {uploaded_file.name} — install python-pptx to extract text]"
        except Exception as e:
            return f"[Could not read PPTX: {e}]"

    # ── IMAGES ────────────────────────────────────────────────
    elif filename.endswith((".png", ".jpg", ".jpeg")):
        return f"[Image uploaded: {uploaded_file.name} — describe what's in this image for revision notes]"

    # ── UNKNOWN ───────────────────────────────────────────────
    else:
        return f"[File uploaded: {uploaded_file.name}]"
